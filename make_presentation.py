from pptx import Presentation
from main import get_presentation_content


# user_fullname = input("Enter your full name: ")
user_fullname = "Alisher Shamuratov"
# topic = input("Enter the topic of your presentation: ")
topic = "Cats"

"""
generated_texts = {
    "user_fullname": "text",
    "context_1": "text",
    "context_2": "text",
    "context_body_1": "text",
    "context_body_1_continue": "text",
    "context_body_2": "text",
    "conclusion_body": "text"
}
"""


# Load your PowerPoint template
presentation = Presentation("template.pptx")

generated_texts = get_presentation_content(topic)


# images_to_add = {
#     2: [("images/image1.png", Inches(1), Inches(1), Inches(3), Inches(2))],  # Add to slide 1
#     # 2: [("image2.jpg", Inches(2), Inches(2), Inches(4), Inches(3))],  # Add to slide 3
# }
#
# # Iterate over slides and insert images
# for slide_index, images in images_to_add.items():
#     slide = presentation.slides[slide_index]
#     for image_path, left, top, width, height in images:
#         slide.shapes.add_picture(image_path, left, top, width, height)

# Iterate through slides and replace text
for slide in presentation.slides:
    for shape in slide.shapes:
        if shape.has_text_frame:
            for paragraph in shape.text_frame.paragraphs:
                for key, replacement in generated_texts.items():
                    if key in paragraph.text.lower():
                        paragraph.text = replacement

# Save the updated presentation
presentation.save("updated_presentation.pptx")
