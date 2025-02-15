import os
import re
from pptx import Presentation
from pptx.util import Inches
import google.generativeai as genai
from pptx.dml.color import RGBColor

TOKEN = 'AIzaSyCXRYaZTzqhdCCtYlya8BW83_VAEyTjR1w'
genai.configure(api_key=TOKEN)

model = genai.GenerativeModel("gemini-1.5-flash")

TEMPLATES = {
    "design0": "design0.pptx",
    "design1": "design1.pptx",
    "design2": "design2.pptx",
    "Custom Template": "custom_template.pptx"
}


def get_presentation_content(query):
    """
    Generates structured text content for a presentation on a given topic.

    Args:
      query (str): The topic for the presentation.

    Returns:
      dict: A dictionary containing generated text for different sections.
    """

    prompt = f"""
    Create a structured PowerPoint presentation on **"{query}"**.
    Use the following clear section labels:

    **Title:** (A suitable title for {query} the presentation)
    **Context 1:** (Introduction to, 3-4 sentences {query})
    **Context 2:** (Additional background, 3-4 sentences {query})
    **History:** (Brief historical background if applicable, else state 'No history available' {query})
    **Key Points:** (Main facts, characteristics, or arguments in bullet points)
    **Conclusion:** (Summary and final thoughts in 3-4 sentences)

    Ensure each section is clearly labeled with "**Section Name:**" followed by the content.
    """

    response = model.generate_content(prompt)

    # Validate response
    if not response.text:
        return {"error": "Failed to generate response"}

    # Process response
    sections = ["Context 1", "Context 2", "History", "Key Points", "Conclusion"]
    generated_texts = {key.lower().replace(" ", "_"): "" for key in sections}

    for key in sections:
        pattern = rf"\*\*\s*Section Name:\s*{re.escape(key)}.*?\*\*\s*(.*?)\n\s*(?=\*\*Slide|\*\*Section Name|\Z)"
        match = re.search(pattern, response.text, re.DOTALL)

        if match:
            extracted_text = match.group(1).strip()
            generated_texts[key.lower().replace(" ", "_")] = extracted_text
            print(f'Matched "{key}":, {extracted_text[:50]}....')
        else:
            print(f'No match for "{key}"')

    return generated_texts


def create_presentation(topic, template_choices="design0.pptx", output_folder="p"):
    template_path = TEMPLATES.get(template_choices)

    if not template_path or not os.path.exists(template_path):
        print("Template not found. Falling back to default template.")
        prs = Presentation()
    else:
        prs = Presentation(template_path)

    content = get_presentation_content(topic)

    # Handle errors
    if "error" in content:
        print(content["error"])
        return

    # Ensure output folder exists
    os.makedirs(output_folder, exist_ok=True)

    # Define PowerPoint file path
    pptx_file = os.path.join(output_folder, f"{topic.replace(' ', '_')}.pptx")

    # Add the title slide FIRST
    slide_layout = prs.slide_layouts[0]  # Title slide layout
    slide = prs.slides[0] if len(prs.slides) > 0 else prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]

    # Use the topic as the title or generate a title from the content if available
    title_text = content.get("title", topic)  # This allows flexibility if "title" is missing
    title.text = title_text
    subtitle.text = "Your Name"

    # Overwrite slides based on the template (excluding the title slide)
    sections = ["context_1", "context_2", "history", "key_points", "conclusion"]

    for idx, section in enumerate(sections):
        if content.get(section):  # Only overwrite slide if there's content
            # Get slide
            slide = prs.slides[idx + 1] if idx + 1 < len(prs.slides) else prs.slides.add_slide(prs.slide_layouts[1])

            # Overwrite slide title and content
            if slide.shapes.title:
                slide.shapes.title.text = section.replace("_", " ").title()

            if len(slide.placeholders) > 1:
                text_box = slide.placeholders[1]  # Ensure this is a content placeholder
                text_frame = text_box.text_frame
                text_frame.clear()  # Clears any default text

                if section == "key_points":
                    for line in content.get(section, "").split("\n"):
                        if line.strip():
                            p = text_frame.add_paragraph()
                            p.text = line.strip()
                            p.level = 0
                else:
                    text_frame.text = content.get(section, "")

    # Save the presentation with overwritten content
    prs.save(pptx_file)
    print(f"✅ Presentation saved at: {pptx_file}")



if __name__ == "__main__":

    print("Choose a template:")
    for idx, template in enumerate(TEMPLATES.keys(), 1):
        print(f"{idx}. {template}")

    choice = input("Enter the number, template choice: ")
    try:
        template_choice = list(TEMPLATES.keys())[int(choice) - 1]
    except (IndexError, ValueError):
        template_choice = "design0"  # Default if the input is invalid

    topic = input("Enter your topic: ")
    create_presentation(topic, template_choices=template_choice)
